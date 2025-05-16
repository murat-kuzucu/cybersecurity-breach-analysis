import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pathlib import Path
import matplotlib.pyplot as plt

# Create a presentation
def create_cybersecurity_presentation():
    prs = Presentation()
    
    # Title slide
    slideLayout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slideLayout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Cybersecurity Breach Analysis"
    subtitle.text = "2015-2024 Global Trends and Patterns"
    
    # Introduction slide
    slideLayout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slideLayout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Research Question"
    content.text = "What are the most significant factors affecting financial losses and resolution times in cybersecurity breaches across different industries and regions?"
    
    # Methodology slide
    slideLayout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slideLayout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Methodology"
    content.text = "• Dataset: Global Cybersecurity Threats 2015-2024\n"
    content.text += "• Analysis of financial impacts across industries\n"
    content.text += "• Evaluation of resolution times by attack type\n"
    content.text += "• Examination of security vulnerability effectiveness\n"
    content.text += "• Analysis of attack source patterns\n"
    content.text += "• Supplementary questionnaire for qualitative insights"
    
    # Main findings 1
    slideLayout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slideLayout)
    title = slide.shapes.title
    
    title.text = "Key Finding 1: Financial Impact by Industry"
    
    # Add picture if it exists
    imagePath = Path('analysis_results/avg_loss_by_industry.png')
    if imagePath.exists():
        left = Inches(1.5)
        top = Inches(2)
        width = Inches(7)
        slide.shapes.add_picture(str(imagePath), left, top, width=width)
    
    # Main findings 2
    slideLayout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slideLayout)
    title = slide.shapes.title
    
    title.text = "Key Finding 2: Financial Loss Trends Over Time"
    
    # Add picture if it exists
    imagePath = Path('analysis_results/loss_trend_over_time.png')
    if imagePath.exists():
        left = Inches(1.5)
        top = Inches(2)
        width = Inches(7)
        slide.shapes.add_picture(str(imagePath), left, top, width=width)
    
    # Resolution time analysis
    slideLayout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slideLayout)
    title = slide.shapes.title
    
    title.text = "Key Finding 3: Attack Resolution Times"
    
    # Add picture if it exists
    imagePath = Path('analysis_results/avg_resolution_by_attack.png')
    if imagePath.exists():
        left = Inches(1.5)
        top = Inches(2)
        width = Inches(7)
        slide.shapes.add_picture(str(imagePath), left, top, width=width)
    
    # Vulnerability analysis
    slideLayout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slideLayout)
    title = slide.shapes.title
    
    title.text = "Key Finding 4: Security Vulnerabilities and Financial Impact"
    
    # Add picture if it exists
    imagePath = Path('analysis_results/loss_by_vulnerability.png')
    if imagePath.exists():
        left = Inches(1.5)
        top = Inches(2)
        width = Inches(7)
        slide.shapes.add_picture(str(imagePath), left, top, width=width)
    
    # Attack sources
    slideLayout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slideLayout)
    title = slide.shapes.title
    
    title.text = "Key Finding 5: Attack Source Distribution"
    
    # Add picture if it exists
    imagePath = Path('analysis_results/attack_source_distribution.png')
    if imagePath.exists():
        left = Inches(1.5)
        top = Inches(2)
        width = Inches(7)
        slide.shapes.add_picture(str(imagePath), left, top, width=width)
    
    # Recommendations slide
    slideLayout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slideLayout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Recommendations"
    content.text = "1. Implement layered security with emphasis on AI-based detection\n\n"
    content.text += "2. Prioritize rapid incident response capabilities\n\n"
    content.text += "3. Deploy industry-specific security strategies\n\n"
    content.text += "4. Enhance security awareness training to address social engineering\n\n"
    content.text += "5. Maintain rigorous software patching practices"
    
    # Conclusion slide
    slideLayout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slideLayout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Conclusion"
    content.text = "• Significant variations in financial impact across industries\n\n"
    content.text += "• Strong correlation between resolution time and financial damage\n\n"
    content.text += "• Social engineering and unpatched software remain top vulnerabilities\n\n"
    content.text += "• AI-based detection systems demonstrate superior effectiveness\n\n"
    content.text += "• Industry-specific security approaches are essential"
    
    # Q&A slide
    slideLayout = prs.slide_layouts[5]  # Title only layout
    slide = prs.slides.add_slide(slideLayout)
    title = slide.shapes.title
    
    title.text = "Questions & Discussion"
    
    # Save the presentation
    prs.save('presentation.pptx')
    print("Presentation created successfully.")

if __name__ == "__main__":
    create_cybersecurity_presentation()
